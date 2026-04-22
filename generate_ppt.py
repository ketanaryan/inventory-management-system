import sys
from pptx import Presentation
from pptx.util import Pt

def set_text(shape, text, is_title=False):
    if not shape.has_text_frame:
        return
    text_frame = shape.text_frame
    text_frame.clear()  # Clear existing content
    p = text_frame.paragraphs[0]
    p.text = text
    if is_title:
        p.font.size = Pt(36)
    else:
        p.font.size = Pt(24)

def add_bullets(shape, bullets):
    if not shape.has_text_frame:
        return
    text_frame = shape.text_frame
    text_frame.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)

def main():
    prs = Presentation('PPT template for Final Presentation 2025-26.pptx')
    
    # Slide 1: Title
    # Shape 0 is the title placeholder. Shape 1 is students.
    set_text(prs.slides[0].shapes[0], "Blockchain-based Inventory Management System", is_title=True)
    # Keeping students as is, or you can manually update
    
    # Slide 2: Content
    content_bullets = [
        "Introduction",
        "Literature Review",
        "Problem Statement",
        "Objectives",
        "Proposed Work",
        "Results Analysis",
        "Conclusion",
        "Future Scope",
        "References"
    ]
    add_bullets(prs.slides[1].shapes[1], content_bullets)
    
    # Slide 3: Introduction
    intro_bullets = [
        "Traditional inventory systems suffer from a lack of transparency and are prone to data manipulation.",
        "Our project leverages blockchain technology to build a decentralized, secure inventory management system.",
        "It provides an immutable ledger where all transactions (adding, updating, tracking items) are permanently recorded.",
        "Ensures trust among stakeholders without relying on a central authority."
    ]
    add_bullets(prs.slides[2].shapes[1], intro_bullets)
    
    # Slide 4: Literature Review
    lit_bullets = [
        "Traditional ERP systems rely on centralized databases, creating single points of failure and security risks.",
        "Recent studies emphasize the use of blockchain for supply chain management to enhance traceability.",
        "Existing blockchain solutions are often costly or overly complex for small-to-medium enterprises.",
        "We propose a lightweight, Ethereum-based Decentralized Application (DApp) to address these gaps."
    ]
    add_bullets(prs.slides[3].shapes[1], lit_bullets)
    
    # Slide 5: Problem Statement
    prob_bullets = [
        "Lack of trust and transparency in multi-party supply chains.",
        "High risk of data tampering or loss in centralized databases.",
        "Difficult to definitively trace the origin and lifecycle of items in conventional systems.",
        "Need for a transparent, secure, and decentralized inventory tracking mechanism."
    ]
    add_bullets(prs.slides[4].shapes[1], prob_bullets)
    
    # Slide 6: Objectives
    obj_bullets = [
        "To develop a decentralized inventory management system using Ethereum smart contracts.",
        "To ensure data immutability and transparency for all product lifecycle events.",
        "To provide a user-friendly frontend (Next.js) for interacting seamlessly with the blockchain.",
        "To reduce fraud and administrative overhead in supply chain management."
    ]
    add_bullets(prs.slides[5].shapes[1], obj_bullets)
    
    # Slide 7: Proposed Work
    # Shape 0 is content textbox
    prop_bullets = [
        "System Architecture: Modern Next.js Frontend styled with Tailwind CSS.",
        "Smart Contracts: Written in Solidity and deployed on an Ethereum network (Ganache/Truffle).",
        "Blockchain Integration: Ethers.js library used for communication between UI and Blockchain.",
        "Core Modules: Add New Item, Update Item Status, Track Item History."
    ]
    try:
        add_bullets(prs.slides[6].shapes[0], prop_bullets)
    except:
        pass
    
    # Slide 8: Results Analysis
    res_bullets = [
        "Smart contracts successfully deployed with optimized gas costs.",
        "Instant, secure, and verifiable transactions across the network.",
        "Immutable transaction history prevents unauthorized data alterations.",
        "Intuitive dashboard interface allows real-time tracking of items for users."
    ]
    add_bullets(prs.slides[7].shapes[2], res_bullets)
    
    # Slide 9: Conclusion
    # Shape 3 is content
    conc_bullets = [
        "The blockchain-based inventory system successfully demonstrates decentralized tracking.",
        "It eliminates the need for a central authority, enhancing trust among all parties.",
        "Data tampering is virtually impossible due to the immutable nature of the ledger.",
        "Web3 technologies provide a robust, modern alternative to traditional centralized ERPs."
    ]
    try:
        add_bullets(prs.slides[8].shapes[3], conc_bullets)
    except:
        pass
    
    # Slide 10: Future Scope
    fut_bullets = [
        "Integration with IoT sensors (RFID, Barcodes) for automated item tracking and updates.",
        "Implementing cross-chain compatibility for broader industry adoption.",
        "Adding advanced analytics and machine learning algorithms for demand forecasting.",
        "Deploying the system to the Ethereum mainnet or scalable Layer 2 solutions (e.g., Polygon)."
    ]
    add_bullets(prs.slides[9].shapes[2], fut_bullets)
    
    # Slide 11: Authors Publication
    pub_bullets = [
        "Currently preparing research paper for submission.",
        "Target: IEEE International Conference on Blockchain Technology.",
        "Focus: \"Decentralized Approaches to Supply Chain Traceability using Smart Contracts.\""
    ]
    add_bullets(prs.slides[10].shapes[2], pub_bullets)
    
    # Slide 12: References
    ref_bullets = [
        "[1] Nakamoto, S. 'Bitcoin: A Peer-to-Peer Electronic Cash System', 2008.",
        "[2] Wood, G. 'Ethereum: A Secure Decentralised Generalised Transaction Ledger', 2014.",
        "[3] Swan, M. 'Blockchain: Blueprint for a New Economy', O'Reilly Media, 2015.",
        "[4] Various research papers on Blockchain in Supply Chain Management and IoT integration."
    ]
    add_bullets(prs.slides[11].shapes[2], ref_bullets)
    
    # Save the new presentation
    output_path = 'Final_Presentation_Inventory_System.pptx'
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    main()
