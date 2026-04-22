import sys
from pptx import Presentation

def analyze_ppt(file_path):
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for j, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            text = shape.text.replace("\n", " ").strip()
            if text:
                print(f"Shape {j} ({shape.shape_type}): {text[:100]}")

if __name__ == "__main__":
    analyze_ppt('PPT template for Final Presentation 2025-26.pptx')
